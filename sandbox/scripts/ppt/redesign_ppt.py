from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_redesign_ppt(input_path, output_path):
    prs = Presentation()

    # --- Color Palette ---
    COLOR_NAVY = RGBColor(0, 32, 96)
    COLOR_BLUE = RGBColor(0, 112, 192)
    COLOR_LIGHT_BLUE = RGBColor(221, 235, 247)
    COLOR_GRAY = RGBColor(128, 128, 128)
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_ACCENT = RGBColor(192, 0, 0) # For warnings/traps
    COLOR_TEXT = RGBColor(32, 32, 32)

    def apply_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def set_text_style(shape, font_size=Pt(18), bold=False, color=COLOR_TEXT, align=PP_ALIGN.LEFT):
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = align
            for run in paragraph.runs:
                run.font.size = font_size
                run.font.bold = bold
                run.font.color.rgb = color
                run.font.name = 'Arial'

    # --- Slide Data (Extracted from previous analysis) ---
    # Note: Reconstructing text accurately from the garbled output is tricky.
    # I will use the structure and key terms identified.
    
    slides_data = [
        {
            "type": "title",
            "title": "クリティカル・シンキング・トレーニング",
            "subtitle": "論理的な思考で、より良い意思決定を\n\n講師: ____________________\n日付: ____________________"
        },
        {
            "type": "content",
            "title": "クリティカル・シンキングとは",
            "content": [
                "定義: 「批判的」ではなく「客観的」に考えること",
                "目的: 情報を鵜呑みにせず、論理的な根拠に基づいて判断すること",
                "重要性:",
                "  ・情報の真偽を見極める力",
                "  ・偏った判断（バイアス）を防ぐ力",
                "  ・より質の高い意思決定を行う力"
            ]
        },
        {
            "type": "content",
            "title": "なぜ必要なのか？",
            "content": [
                "「思い込み」によるミスを防ぐ",
                "例: 「いつもこうだから」という経験則への過度な依存",
                "「論理的な思考」への転換:",
                "  ・根拠（Fact）と意見（Opinion）を区別する"
            ]
        },
        {
            "type": "content",
            "title": "クリティカル・シンキングのメリット",
            "content": [
                "判断の精度向上: 根拠に基づいた確実な判断",
                "問題解決力の向上: 本質的な原因の特定",
                "コミュニケーションの改善: 論理的な説明による説得力"
            ]
        },
        {
            "type": "section",
            "title": "思考の3ステップ",
            "content": ["Step 1: 情報の整理", "Step 2: 批判的検討", "Step 3: 事実と意見の分離"]
        },
        {
            "type": "content",
            "title": "Step 1: 情報の整理",
            "content": [
                "情報を整理し、全体像を把握する",
                "ポイント:",
                "  ・情報の出所（ソース）を確認する",
                "  ・情報の鮮度（いつの情報か）を確認する",
                "  ・情報の偏り（誰が言っているか）を確認する"
            ]
        },
        {
            "type": "content",
            "title": "Step 2: 批判的検討",
            "content": [
                "「本当にそうか？」と問い直す",
                "検討項目:",
                "  ・前提条件は正しいか？",
                "  ・論理の飛躍はないか？"
            ]
        },
        {
            "type": "content",
            "title": "Step 3: 事実と意見の分離",
            "content": [
                "Fact（事実）と Opinion（意見）を明確に区別する",
                "チェックリスト:",
                "  ・それは客観的に証明できることか？",
                "  ・それは個人の感想や推測ではないか？",
                "  ・それは論理的な結論か？"
            ]
        },
        {
            "type": "content",
            "title": "陥りやすい罠",
            "content": [
                "「思い込み」の罠:",
                "  ・自分の経験や知識に固執してしまうこと",
                "「バイアス」の罠:",
                "  ・自分に都合の良い情報だけを集めてしまうこと"
            ],
            "is_warning": True
        },
        {
            "type": "content",
            "title": "バイアスの例",
            "content": [
                "確証バイアス: 自分の考えを裏付ける情報ばかり集める",
                "利用可能性ヒューリスティック: 思い出しやすい情報を優先する"
            ]
        },
        {
            "type": "content",
            "title": "思考を深めるヒント",
            "content": [
                "「なぜ？」を繰り返す",
                "「もし〜だったら？」と仮定する",
                "「他の視点はないか？」と考える"
            ]
        },
        {
            "type": "content",
            "title": "実践的なトレーニング",
            "content": [
                "ケーススタディ: 実際の事例を用いて考える",
                "ワークショップ: チームで異なる視点を出し合う"
            ]
        },
        {
            "type": "content",
            "title": "まとめ",
            "content": [
                "クリティカル・シンキングは「技術」である",
                "日々の習慣として、常に問い続けること",
                "より良い判断が、より良い未来を作る"
            ]
        }
    ]

    for data in slides_data:
        if data["type"] == "title":
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            apply_slide_background(slide, COLOR_NAVY)
            
            title = slide.shapes.title
            title.text = data["title"]
            set_text_style(title, font_size=Pt(44), bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
            
            subtitle = slide.placeholders[1]
            subtitle.text = data["subtitle"]
            set_text_style(subtitle, font_size=Pt(24), color=COLOR_WHITE, align=PP_ALIGN.CENTER)

        elif data["type"] == "section":
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            apply_slide_background(slide, COLOR_BLUE)
            
            title = slide.shapes.title
            title.text = data["title"]
            set_text_style(title, font_size=Pt(40), bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
            
            # Add content centered
            body = slide.placeholders[1]
            body.text = "\n".join(data["content"])
            set_text_style(body, font_size=Pt(32), color=COLOR_WHITE, align=PP_ALIGN.CENTER)

        elif data["type"] == "content":
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            
            is_warning = data.get("is_warning", False)
            bg_color = COLOR_LIGHT_BLUE if not is_warning else RGBColor(255, 235, 235)
            apply_slide_background(slide, bg_color)

            title = slide.shapes.title
            title.text = data["title"]
            title_color = COLOR_NAVY if not is_warning else COLOR_ACCENT
            set_text_style(title, font_size=Pt(36), bold=True, color=title_color)

            body = slide.placeholders[1]
            body.text = ""
            for line in data["content"]:
                p = body.text_frame.add_paragraph()
                p.text = line
                # Simple indentation logic
                if line.strip().startswith("  "):
                    p.level = 1
                elif line.strip().startswith("    "):
                    p.level = 2
                else:
                    p.level = 0
            
            text_color = COLOR_TEXT if not is_warning else COLOR_ACCENT
            set_text_style(body, font_size=Pt(22), color=text_color)

    prs.save(output_path)
    print(f"Successfully saved to {output_path}")

if __name__ == "__main__":
    import sys
    input_file = 'output/docs/ppt/Critical_Thinking_Training.pptx'
    output_file = 'output/docs/ppt/Critical_Thinking_Training_Redesign.pptx'
    create_redesign_ppt(input_file, output_file)
