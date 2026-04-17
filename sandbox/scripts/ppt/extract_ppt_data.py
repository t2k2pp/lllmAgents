from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def extract_data(input_path):
    prs = Presentation(input_path)
    extracted_slides = []
    
    for slide in prs.slides:
        slide_data = {
            "layout_name": slide.slide_layout.name,
            "elements": []
        }
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    # shapeのタイプや位置などの情報を保持
                    element = {
                        "text": text,
                        "type": "text",
                        "top": shape.top,
                        "left": shape.left,
                        "width": shape.width,
                        "height": shape.height,
                        "paragraphs": []
                    }
                    
                    # 各段落の詳細（インデント、フォントサイズ等）を抽出
                    for paragraph in shape.text_frame.paragraphs:
                        para_data = {
                            "text": paragraph.text,
                            "level": paragraph.level,
                            "alignment": paragraph.alignment,
                            "runs": []
                        }
                        for run in paragraph.runs:
                            para_data["runs"].append({
                                "text": run.text,
                                "font_size": run.font.size,
                                "bold": run.font.bold,
                                "color": run.font.color.rgb if run.font.color and hasattr(run.font.color, 'rgb') else None
                            })
                        element["paragraphs"].append(para_data)
                    
                    slide_data["elements"].append(element)
        
        extracted_slides.append(slide_data)
    return extracted_slides

if __name__ == "__main__":
    import json
    input_file = 'output/docs/ppt/Critical_Thinking_Training.pptx'
    data = extract_data(input_file)
    # JSONとして保存して後続のプロセスで利用可能にする
    with open('extracted_ppt_data.json', 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    print(f"Successfully extracted data from {input_file} to extracted_ppt_data.json")
