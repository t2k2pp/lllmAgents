import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def redesign_ppt(input_json_path, input_pptx_path, output_pptx_path):
    with open(input_json_path, 'r', encoding='utf-8') as f:
        extracted_slides = json.load(f)

    prs = Presentation(input_pptx_path)
    
    # Color Palette
    COLOR_NAVY = RGBColor(0, 32, 96)
    COLOR_BLUE = RGBColor(0, 112, 192)
    COLOR_LIGHT_BLUE = RGBColor(221, 235, 247)
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_TEXT = RGBColor(32, 32, 32)
    COLOR_ACCENT = RGBColor(192, 0, 0)

    for i, (original_slide, data) in enumerate(zip(prs.slides, extracted_slides)):
        # 1. Background Color
        background = original_slide.background
        fill = background.fill
        fill.solid()
        
        if data["layout_name"] == "Title Slide":
            fill.fore_color.rgb = COLOR_NAVY
        elif any(keyword in str(data["elements"]) for keyword in ["Step", "フロー", "アクション"]):
            fill.fore_color.rgb = COLOR_BLUE
        else:
            fill.fore_color.rgb = COLOR_LIGHT_BLUE

        # 2. Update all text shapes
        # We iterate through the extracted elements and find the best matching shape in the original slide
        # by position (top, left) to preserve layout.
        
        for element in data["elements"]:
            # Find the shape in the original slide that is closest to this element's position
            best_shape = None
            min_dist = float('inf')
            
            target_top = element["top"]
            target_left = element["left"]

            for shape in original_slide.shapes:
                if shape.has_text_frame:
                    # Calculate Euclidean distance between centers (roughly)
                    dist = ((shape.top - target_top)**2 + (shape.left - target_left)**2)**0.5
                    if dist < min_dist:
                        min_dist = dist
                        best_shape = shape
            
            if best_shape and min_dist < 5000000: # Threshold to ensure we match reasonably well
                tf = best_shape.text_frame
                tf.clear() # Clear existing text
                
                for p_idx, p_data in enumerate(element["paragraphs"]):
                    # Add paragraph
                    if p_idx == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    
                    p.text = p_data["text"]
                    p.level = p_data["level"]
                    
                    # Alignment
                    if p_data["alignment"]:
                        p.alignment = p_data["alignment"]
                    elif data["layout_name"] == "Title Slide":
                        p.alignment = PP_ALIGN.CENTER
                    
                    # Font Style
                    for run in p.runs:
                        run.font.name = 'Arial'
                        
                        # Base color and size
                        if data["layout_name"] == "Title Slide":
                            run.font.color.rgb = COLOR_WHITE
                            run.font.size = Pt(32) if p_data["level"] == 0 else Pt(20)
                        else:
                            run.font.color.rgb = COLOR_TEXT
                            run.font.size = Pt(24) if p_data["level"] == 0 else Pt(18)
                        
                        # Boldness
                        if p_data["level"] == 0:
                            run.font.bold = True
                        else:
                            run.font.bold = False
                            
                        # Accent color for titles in content slides
                        if data["layout_name"] != "Title Slide" and p_data["level"] == 0:
                            # If it's at the top of the slide, make it navy
                            if element["top"] < 1000000:
                                run.font.color.rgb = COLOR_NAVY
                                run.font.size = Pt(32)

    prs.save(output_pptx_path)
    print(f"Redesign completed: {output_pptx_path}")

if __name__ == "__main__":
    import sys
    input_json = 'extracted_ppt_data.json'
    input_pptx = 'output/docs/ppt/Critical_Thinking_Training.pptx'
    output_pptx = 'output/docs/ppt/Critical_Thinking_Training_Redesign_V3.pptx'
    redesign_ppt(input_json, input_pptx, output_pptx)
