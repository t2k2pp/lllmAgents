import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def redesign_ppt(input_json_path, input_pptx_path, output_pptx_path):
    with open(input_json_path, 'r', encoding='utf-8') as f:
        extracted_slides = json.load(f)

    prs = Presentation(input_pptx_path)
    
    # Color Palette
    COLOR_NAVY = RGBColor(0, 32, 96)
    COLOR_BLUE = RGBColor(0, 112, 192)
    COLOR_LIGHT_BLUE = RGBColor(221, 235, 247)
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_TEXT = RGBColor(32, 32, 32)
    COLOR_ACCENT = RGBColor(192, 0, 0)

    def apply_style(slide, layout_name, elements):
        # Background color based on layout
        background = slide.background
        fill = background.fill
        fill.solid()
        
        if layout_name == "Title Slide":
            fill.fore_color.rgb = COLOR_NAVY
        elif "Step" in str(elements) or "フロー" in str(elements): # Heuristic for section/process
             fill.fore_color.rgb = COLOR_BLUE
        else:
            fill.fore_color.rgb = COLOR_LIGHT_BLUE

        for element in elements:
            if element["type"] == "text":
                # Find or create shape
                # In a real redesign, we'd use the original position but adjust for new styles
                # For this task, we'll use the extracted position but ensure it's a text shape
                
                # To avoid creating duplicate shapes in existing slides, 
                # we actually should be clearing old ones or using the existing ones.
                # But since we are 'redesigning', we'll create a new slide with the same content.
                pass

    # Since we want to REUSE the existing slide objects but CHANGE their appearance:
    # We will iterate through the original slides and the extracted data in parallel.
    
    for i, (original_slide, data) in enumerate(zip(prs.slides, extracted_slides)):
        # 1. Apply Background
        background = original_slide.background
        fill = background.fill
        fill.solid()
        
        if data["layout_name"] == "Title Slide":
            fill.fore_color.rgb = COLOR_NAVY
        else:
            fill.fore_color.rgb = COLOR_LIGHT_BLUE

        # 2. Update all text shapes in the original slide
        # We'll use the extracted data to "reset" the text and apply new styles
        # to ensure we don't miss anything and keep it clean.
        
        # First, clear existing text frames to avoid overlap if we add new ones
        # However, it's safer to just iterate through existing shapes and update them.
        
        # For simplicity and robustness in this agentic context, 
        # we will map the extracted paragraphs back to the existing shapes.
        
        # Sort elements by top position to match them to existing shapes
        sorted_elements = sorted(data["elements"], key=lambda x: x["top"])
        existing_shapes = sorted(original_slide.shapes, key=lambda x: x.top)
        
        # Match shapes to elements (this is a heuristic)
        for j, element in enumerate(sorted_elements):
            if j < len(existing_shapes):
                shape = existing_shapes[j]
                if shape.has_text_frame:
                    tf = shape.text_frame
                    tf.clear() # Clear existing
                    
                    for p_idx, p_data in enumerate(element["paragraphs"]):
                        p = tf.paragraphs[0] if p_idx == 0 else tf.add_paragraph()
                        p.text = p_data["text"]
                        p.level = p_data["level"]
                        
                        # Alignment
                        if p_data["alignment"]:
                            p.alignment = p_data["alignment"]
                        elif data["layout_name"] == "Title Slide":
                            p.alignment = PP_ALIGN.CENTER
                        
                        # Font Style
                        for run in p.runs:
                            run.font.name = 'Arial'
                            run.font.size = Pt(24) if data["layout_name"] == "Title Slide" else Pt(18)
                            run.font.bold = True if p_data["level"] == 0 else False
                            
                            # Color
                            if data["layout_name"] == "Title Slide":
                                run.font.color.rgb = COLOR_WHITE
                            else:
                                run.font.color.rgb = COLOR_TEXT
                                
                                # If it's a title-like element (top of slide)
                                if element["top"] < 1000000:
                                    run.font.size = Pt(32)
                                    run.font.bold = True
                                    run.font.color.rgb = COLOR_NAVY

    prs.save(output_pptx_path)
    print(f"Redesign completed: {output_pptx_path}")

if __name__ == "__main__":
    import sys
    input_json = 'extracted_ppt_data.json'
    input_pptx = 'output/docs/ppt/Critical_Thinking_Training.pptx'
    output_pptx = 'output/docs/ppt/Critical_Thinking_Training_Redesign_V2.pptx'
    redesign_ppt(input_json, input_pptx, output_pptx)
